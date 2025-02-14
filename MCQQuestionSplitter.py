import sys
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
import os
import tempfile
import re
from PIL import Image
from tqdm import tqdm

class TemplateManager:
    @staticmethod
    def get_template_path():
        """Get the path to the PowerPoint template file."""
        if getattr(sys, 'frozen', False):
            # Running as compiled executable
            base_path = sys._MEIPASS
        else:
            # Running as script
            base_path = os.path.dirname(os.path.abspath(__file__))
        
        template_path = os.path.join(base_path, 'templates', 'default.pptx')
        
        # If template doesn't exist in the expected location, extract it
        if not os.path.exists(template_path):
            TemplateManager.extract_template(template_path)
            
        return template_path
    
    @staticmethod
    def extract_template(target_path):
        """Extract the default PowerPoint template to the specified location."""
        # Create templates directory if it doesn't exist
        os.makedirs(os.path.dirname(target_path), exist_ok=True)
        
        # Create a blank presentation and save it as the template
        prs = Presentation()
        prs.save(target_path)

class MCQQuestionSplitter:
    def __init__(self, slide_duration=None):
        self.temp_dir = tempfile.mkdtemp()
        self.slide_duration = slide_duration  # Can be None for manual slide control
        self.template_path = TemplateManager.get_template_path()

    def detect_questions(self, pdf_path):
        """Detect questions with consistent formatting and clear boundaries."""
        questions = []
        current_question = None
        expected_question = 1
        reference_formatting = None
        
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages[1:], 1):  # Start from page 0
                # Extract words with their properties
                words = page.extract_words(
                    keep_blank_chars=True,
                    x_tolerance=3,
                    y_tolerance=3,
                    extra_attrs=['fontname', 'size', 'object_type']
                )
                
                # Group words into lines based on vertical position
                current_y = None
                current_line = []
                lines = []
                
                for word in words:
                    if current_y is None:
                        current_y = word['top']
                        current_line.append(word)
                    elif abs(word['top'] - current_y) <= 3:  # Same line
                        current_line.append(word)
                    else:  # New line
                        lines.append(current_line)
                        current_line = [word]
                        current_y = word['top']
                
                if current_line:
                    lines.append(current_line)
                                
                # Process each line
                for line in lines:
                    # Combine words into line text
                    line_text = ' '.join(word['text'] for word in line)
                    
                    # Check for question patterns
                    num_match = re.match(r'^\s*(\d+)[.\s]', line_text)

                    # If first question, set reference formatting
                    if not reference_formatting and num_match:
                        reference_formatting = {
                            'fontname': line[0].get('fontname'),
                            'size': line[0].get('size'),
                            'color': line[0].get('strokedColor')
                        }
                    
                    # Validate formatting for subsequent questions
                    formatting_match = (
                        reference_formatting and 
                        line[0].get('fontname') == reference_formatting['fontname'] and
                        abs(line[0].get('size', 0) - reference_formatting['size']) <= 1
                    )
                    
                    question_starters = r'(Which|What|How|Why|Where|When|Whose|Who|In|The|If|Define|State|Calculate)'
                    starter_match = re.match(f"^{question_starters}", line_text, re.IGNORECASE)
                                        
                    is_question = False
                    if num_match:
                        question_num = int(num_match.group(1))
                        is_question = question_num == expected_question and formatting_match
                    elif starter_match and len(line_text.split()) > 3:
                        is_question = formatting_match and not current_question
                    
                    if is_question:
                        # Finalize previous question
                        if current_question:
                            questions.append(current_question)
                        
                        bbox = [line[0]['x0'], line[0]['top'], 
                            line[-1]['x1'], line[-1]['bottom']]
                        
                        current_question = {
                            'number': expected_question,
                            'page': page_num,
                            'start_bbox': bbox,
                            'end_bbox': bbox.copy(),
                            'content': [(page_num, bbox, line_text)]
                        }
                        options_cnt = 0
                        expected_question += 1
                        print(f"Found question {expected_question-1}: {line_text}")
                    
                    elif current_question and line_text != " " and options_cnt < 4:
                        # Capture all content between questions
                        bbox = [line[0]['x0'], line[0]['top'], 
                            line[-1]['x1'], line[-1]['bottom']]
                        # print(line_text, bbox, options_cnt)
                        current_question['content'].append((page_num, bbox, line_text))
                        current_question['end_bbox'][2] = max(current_question['end_bbox'][2], bbox[2])
                        current_question['end_bbox'][3] = max(current_question['end_bbox'][3], bbox[3])

                    if re.match(r'[A-D]\s+[^)]+\s+[A-D]\s+[^)]+\s+[A-D]\s+[^)]+\s+[A-D]\s+[^)]+', line_text) or line_text == 'A B C D':
                        options_cnt += 4
                    elif re.match(r'^[A-D]\s+[A-D]', line_text):
                        options_cnt += 2
                    elif re.search(r'[A-D]\s+.+', line_text) or line_text in 'ABCD':
                        options_cnt += 1
                        
            # Add the last question
            if current_question:
                questions.append(current_question)
                

        return questions

    def capture_question_image(self, pdf_path, question, questions):
        """Capture entire question including images up until the next question starts."""
        with pdfplumber.open(pdf_path) as pdf:
            page = pdf.pages[question['page']]
            
            # Calculate initial boundary from current question
            bbox = question['start_bbox'].copy()
            
            # Find the next question that appears after this one
            next_question = None
            for q in questions:
                if q['number'] == question['number'] + 1:
                    next_question = q
                    break
            
            # Update bbox based on content and next question position
            for content in question['content']:
                page_num, content_bbox, _ = content
                
                # If content is on a different page than the next question
                # or if content appears before the next question on the same page
                should_include = True
                if next_question and page_num == next_question['page']:
                    if content_bbox[1] >= next_question['start_bbox'][1]:
                        should_include = False
                
                if should_include:
                    bbox[0] = min(bbox[0], content_bbox[0])
                    bbox[1] = min(bbox[1], content_bbox[1])
                    bbox[2] = max(bbox[2], content_bbox[2])
                    bbox[3] = max(bbox[3], content_bbox[3])
                    # print(bbox)
            
            # If there's a next question on the same page, use its start position
            # as the end boundary
            if next_question and next_question['page'] == question['page']:
                bbox[3] = next_question['start_bbox'][1] - 5  # Small gap
                # print(next_question['start_bbox'][1], bbox[3])
            elif question['page']+1 == len(pdf.pages):
                bbox[3] = bbox[3] + 30
            else:
                # If this is the last question on the page, extend to bottom
                # or if question continues to next page, extend to page bottom
                bbox[3] = max(page.bbox[3] - 50, bbox[3])
            
            # Ensure reasonable width
            bbox[2] = min(bbox[2] * 1.20, page.bbox[2])
            
            # Handle multi-page questions
            # if next_question and next_question['page'] > question['page']:
            #     # Capture full remaining page height for current page
            #     bbox[3] = page.bbox[3]
            
            # Render page to image
            img = page.crop(bbox).to_image(resolution=200)
            
            # Save to temporary file
            img_path = os.path.join(self.temp_dir, f'question_{question["number"]}.png')
            img.save(img_path)
            return img_path

    def set_slide_timing(self, slide, seconds):
        """Set the slide transition to advance automatically after the specified number of seconds."""
        # Skip timing if no duration specified
        if seconds is None:
            return
            
        # Get the slide's XML element
        slide_element = slide._element
        
        # Create transition element if it doesn't exist
        transition = slide_element.find('.//p:transition', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if transition is None:
            transition = OxmlElement('p:transition')
            slide_element.insert(2, transition)
        
        # Set advance timing
        timing = slide_element.find('.//p:timing', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if timing is None:
            timing = OxmlElement('p:timing')
            slide_element.insert(3, timing)
        
        # Create or update advance node
        advance = timing.find('.//p:tn:nodeType', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if advance is None:
            advance = OxmlElement('p:tn')
            timing.append(advance)
        
        # Set the advance after time in the transition
        transition.set('advAuto', '1')
        transition.set('advTm', str(seconds * 1000))  # Convert seconds to milliseconds
        
        # Ensure automatic advancement is enabled
        if hasattr(slide, 'slide_time'):
            slide.slide_time = seconds * 1000

    def create_slide_with_question(self, prs, img_path, question_number):
        """Create a slide with the question image."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add question number
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title_box.text_frame.text = f"Question {question_number}"
        title_box.text_frame.paragraphs[0].font.size = Pt(24)
        title_box.text_frame.paragraphs[0].font.bold = True
        
        # Add image
        img = Image.open(img_path)
        aspect_ratio = img.width / img.height

        # Calculate dimensions to fit slide
        max_width = Inches(9)
        max_height = Inches(6.5)
        
        if aspect_ratio > max_width / max_height:
            width = max_width
            height = width / aspect_ratio
        else:
            height = max_height
            width = height * aspect_ratio
        
        # Position image
        left = Inches(0.4)
        top = Inches(0.8)
        
        slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        return slide

    def convert_pdf_to_slides(self, pdf_path, output_filename="mcq_presentation.pptx"):
        """Convert PDF MCQ paper to PowerPoint presentation with individual questions."""        
        try:
            # Create presentation
            prs = Presentation(self.template_path)
            
            # Add title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = "Multiple Choice Questions"
            if hasattr(title_slide.shapes, 'placeholders') and len(title_slide.shapes.placeholders) > 1:
                if output_filename == 'mcq_presentation.pptx':
                    output_filename = os.path.basename(pdf_path)[:-4] + '_mcq.pptx'
                title_slide.shapes.placeholders[1].text = os.path.basename(pdf_path)

            # Detect questions
            questions = self.detect_questions(pdf_path)
            
            # Process each question
            for question in tqdm(questions, desc='Processing questions', unit='q'):
                try:
                    img_path = self.capture_question_image(pdf_path, question, questions)
                    slide = self.create_slide_with_question(prs, img_path, question['number'])
                    self.set_slide_timing(slide, self.slide_duration)
                except Exception as e:
                    print(f"Error processing question {question['number']}: {str(e)}")
            
            # Save presentation
            prs.save(output_filename)
            print(f"Presentation saved as {output_filename}")
            
        finally:
            # Cleanup temporary files
            self.cleanup()

    def cleanup(self):
        """Clean up temporary files."""
        for file in os.listdir(self.temp_dir):
            try:
                os.remove(os.path.join(self.temp_dir, file))
            except:
                pass
        try:
            os.rmdir(self.temp_dir)
        except:
            pass

def main():
    import argparse
    parser = argparse.ArgumentParser(description='Convert PDF MCQ paper to PowerPoint presentation')
    parser.add_argument('pdf_path', help='Path to the PDF file')
    parser.add_argument('--output', '-o', default='mcq_presentation.pptx',
                      help='Output PowerPoint file name (default: mcq_presentation.pptx)')
    parser.add_argument('--seconds', '-s', type=int, default=None,
                      help='Number of seconds each slide should display (default: None for manual control)')
    
    args = parser.parse_args()
    
    converter = MCQQuestionSplitter(slide_duration=args.seconds)
    converter.convert_pdf_to_slides(args.pdf_path, args.output)

if __name__ == "__main__":
    main()